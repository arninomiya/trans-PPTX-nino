import os
import json
import time
import uuid
import tempfile
import logging
from typing import Dict, List, Tuple

import azure.functions as func
from azure.functions import AuthLevel
import requests
from tenacity import retry, stop_after_attempt, wait_exponential_jitter

from pptx import Presentation
from azure.storage.blob import BlobClient, ContentSettings

from datetime import datetime, timedelta, timezone
from azure.storage.blob import generate_blob_sas, BlobSasPermissions

app = func.FunctionApp(http_auth_level=AuthLevel.FUNCTION)

# ----------------------------
# ENV (local.settings.json / App Settings に入れる)
# ----------------------------
HTTP_TIMEOUT = int(os.getenv("HTTP_TIMEOUT_SECONDS", "60"))
HTTP_CHUNK_SIZE = int(os.getenv("HTTP_CHUNK_SIZE_BYTES", str(1024 * 1024)))  # 1MB

MAX_PPTX_BYTES = int(os.getenv("MAX_PPTX_BYTES", str(110 * 1024 * 1024)))  # 110MBまで
MAX_TEXT_ITEMS = int(os.getenv("MAX_TEXT_ITEMS", "20000"))
MAX_TEXT_LENGTH_PER_ITEM = int(os.getenv("MAX_TEXT_LENGTH_PER_ITEM", "4000"))
TRANSLATOR_BATCH_SIZE = int(os.getenv("TRANSLATOR_BATCH_SIZE", "50"))

TRANSLATOR_ENDPOINT = os.getenv("TRANSLATOR_ENDPOINT", "https://api.cognitive.microsofttranslator.com")
TRANSLATOR_KEY = os.getenv("TRANSLATOR_KEY", "")
TRANSLATOR_REGION = os.getenv("TRANSLATOR_REGION", "")
TRANSLATOR_API_VERSION = os.getenv("TRANSLATOR_API_VERSION", "3.0")

OUTPUT_BLOB_CONNECTION_STRING = os.getenv("OUTPUT_BLOB_CONNECTION_STRING", "")
OUTPUT_BLOB_CONTAINER = os.getenv("OUTPUT_BLOB_CONTAINER", "pptx-out")
OUTPUT_BLOB_CONTENT_TYPE = "application/vnd.openxmlformats-officedocument.presentationml.presentation"


def _json(status: int, payload: dict) -> func.HttpResponse:
    return func.HttpResponse(
        json.dumps(payload, ensure_ascii=False),
        status_code=status,
        mimetype="application/json"
    )


def _now_ms() -> int:
    return int(time.time() * 1000)


@retry(stop=stop_after_attempt(4), wait=wait_exponential_jitter(initial=1, max=10))
def _http_get_stream(url: str) -> requests.Response:
    r = requests.get(url, stream=True, timeout=HTTP_TIMEOUT, allow_redirects=True)
    r.raise_for_status()
    return r


def download_to_tempfile(url: str, logger: logging.Logger) -> Tuple[str, int]:
    """URLからPPTXをストリーミングDLして一時ファイル保存（メモリに載せない）"""
    start = _now_ms()
    r = _http_get_stream(url)

    # 可能なら事前にサイズチェック
    cl = r.headers.get("Content-Length")
    if cl:
        try:
            if int(cl) > MAX_PPTX_BYTES:
                raise ValueError(f"File too large (Content-Length={cl})")
        except ValueError:
            pass

    fd, path = tempfile.mkstemp(suffix=".pptx")
    os.close(fd)

    total = 0
    try:
        with open(path, "wb") as f:
            for chunk in r.iter_content(chunk_size=HTTP_CHUNK_SIZE):
                if not chunk:
                    continue
                f.write(chunk)
                total += len(chunk)
                if total > MAX_PPTX_BYTES:
                    raise ValueError(f"File too large while streaming ({total} bytes)")
    except Exception:
        try:
            os.remove(path)
        except Exception:
            pass
        raise

    logger.info(json.dumps({"event": "download_complete", "bytes": total, "ms": _now_ms() - start}))
    return path, total


def _iter_runs(prs: Presentation):
    """スライド内テキスト（shape + table）を run 単位で列挙"""
    for s_idx, slide in enumerate(prs.slides):
        for sh_idx, shape in enumerate(slide.shapes):

            # テーブル
            if getattr(shape, "has_table", False):
                tbl = shape.table
                for r_idx, row in enumerate(tbl.rows):
                    for c_idx, cell in enumerate(row.cells):
                        tf = cell.text_frame
                        if not tf:
                            continue
                        for p_idx, p in enumerate(tf.paragraphs):
                            for run_idx, run in enumerate(p.runs):
                                yield (s_idx, sh_idx, f"table:{r_idx},{c_idx}", p_idx, run_idx, run)
                continue

            # 通常テキスト
            if getattr(shape, "has_text_frame", False) and shape.text_frame is not None:
                for p_idx, p in enumerate(shape.text_frame.paragraphs):
                    for run_idx, run in enumerate(p.runs):
                        yield (s_idx, sh_idx, "shape", p_idx, run_idx, run)


def extract_items(prs: Presentation, logger: logging.Logger) -> List[Tuple[Tuple, str]]:
    """翻訳対象テキストを抽出（keyで後で戻せる）"""
    items: List[Tuple[Tuple, str]] = []
    for entry in _iter_runs(prs):
        *k, run = entry
        text = (run.text or "").strip()
        if not text:
            continue
        if len(text) > MAX_TEXT_LENGTH_PER_ITEM:
            logger.warning(json.dumps({"event": "skip_long_text", "len": len(text)}))
            continue
        items.append((tuple(k), text))
        if len(items) >= MAX_TEXT_ITEMS:
            logger.warning(json.dumps({"event": "hit_max_text_items", "max": MAX_TEXT_ITEMS}))
            break
    logger.info(json.dumps({"event": "extract_complete", "count": len(items)}))
    return items


def _dedupe(items: List[Tuple[Tuple, str]]) -> Tuple[List[str], Dict[int, List[int]]]:
    """重複除去して翻訳コスト削減"""
    seen: Dict[str, int] = {}
    uniques: List[str] = []
    u2i: Dict[int, List[int]] = {}
    for i, (_, t) in enumerate(items):
        if t in seen:
            u2i[seen[t]].append(i)
        else:
            u = len(uniques)
            seen[t] = u
            uniques.append(t)
            u2i[u] = [i]
    return uniques, u2i


@retry(stop=stop_after_attempt(4), wait=wait_exponential_jitter(initial=1, max=10))
def translator_batch(texts: List[str], to_lang: str, logger: logging.Logger) -> List[str]:
    """
    Translator batch with 429 handling (Retry-After) + gentle throttling.
    """
    if not TRANSLATOR_KEY:
        raise RuntimeError("TRANSLATOR_KEY is empty")

    url = f"{TRANSLATOR_ENDPOINT}/translate?api-version={TRANSLATOR_API_VERSION}&to={to_lang}"
    headers = {
        "Ocp-Apim-Subscription-Key": TRANSLATOR_KEY,
        "Content-Type": "application/json",
        "X-ClientTraceId": str(uuid.uuid4()),
    }
    if TRANSLATOR_REGION:
        headers["Ocp-Apim-Subscription-Region"] = TRANSLATOR_REGION

    body = [{"text": t} for t in texts]

    # Manual retry loop so we can respect Retry-After precisely
    max_attempts = 8
    for attempt in range(1, max_attempts + 1):
        r = requests.post(url, headers=headers, json=body, timeout=HTTP_TIMEOUT)

        if r.status_code == 429:
            # Respect Retry-After if provided; otherwise backoff
            ra = r.headers.get("Retry-After")
            sleep_s = float(ra) if ra and ra.isdigit() else min(2 ** attempt, 30)
            logger.warning(json.dumps({
                "event": "translator_throttled",
                "status": 429,
                "attempt": attempt,
                "sleep_seconds": sleep_s
            }))
            time.sleep(sleep_s)
            continue

        # For other errors, raise
        r.raise_for_status()
        data = r.json()

        out: List[str] = []
        for item in data:
            trans = item.get("translations", [])
            out.append(trans[0]["text"] if trans else "")
        return out

    # If we exhausted retries
    raise RuntimeError("Translator throttling persisted (429) after retries.")

def translate_items(items: List[Tuple[Tuple, str]], to_lang: str, logger: logging.Logger) -> List[str]:
    uniques, u2i = _dedupe(items)
    logger.info(json.dumps({"event": "dedupe", "original": len(items), "unique": len(uniques)}))

    uniq_out = [""] * len(uniques)
    for i in range(0, len(uniques), TRANSLATOR_BATCH_SIZE):
        batch = uniques[i:i + TRANSLATOR_BATCH_SIZE]
        res = translator_batch(batch, to_lang, logger)
        uniq_out[i:i + TRANSLATOR_BATCH_SIZE] = res
                # gentle throttle between batches (helps avoid 429 on F0/S0)
        time.sleep(float(os.getenv("TRANSLATOR_SLEEP_BETWEEN_BATCHES", "0.6")))

    out = [""] * len(items)
    for u, idxs in u2i.items():
        for idx in idxs:
            out[idx] = uniq_out[u]
    return out


def apply_translations(prs: Presentation, items: List[Tuple[Tuple, str]], translated: List[str], logger: logging.Logger):
    key2trg = {k: t for (k, _), t in zip(items, translated)}
    for entry in _iter_runs(prs):
        *k, run = entry
        k = tuple(k)
        if k in key2trg:
            run.text = key2trg[k]
    logger.info(json.dumps({"event": "apply_complete", "count": len(items)}))


def save_to_temp(prs: Presentation, logger: logging.Logger) -> str:
    fd, p = tempfile.mkstemp(suffix=".pptx")
    os.close(fd)
    prs.save(p)
    logger.info(json.dumps({"event": "save_complete", "path": p}))
    return p


def upload_blob(file_path: str, logger: logging.Logger) -> str:
    if not OUTPUT_BLOB_CONNECTION_STRING:
        raise RuntimeError("OUTPUT_BLOB_CONNECTION_STRING is empty")

    blob_name = f"translated/{uuid.uuid4().hex}.pptx"
    blob = BlobClient.from_connection_string(
        OUTPUT_BLOB_CONNECTION_STRING,
        container_name=OUTPUT_BLOB_CONTAINER,
        blob_name=blob_name
    )

    with open(file_path, "rb") as f:
        blob.upload_blob(
            f,
            overwrite=True,
            content_settings=ContentSettings(content_type=OUTPUT_BLOB_CONTENT_TYPE)
        )

    # --- SAS を生成して「誰でもDLできるURL」を返す ---
    # 接続文字列から AccountName / AccountKey を取り出す（簡易パース）
    parts = dict(
        p.split("=", 1) for p in OUTPUT_BLOB_CONNECTION_STRING.split(";") if "=" in p
    )
    account_name = parts.get("AccountName")
    account_key = parts.get("AccountKey")
    if not account_name or not account_key:
        raise RuntimeError("Cannot parse AccountName/AccountKey from connection string")

    sas_hours = int(os.getenv("OUTPUT_SAS_HOURS", "24"))
    sas = generate_blob_sas(
        account_name=account_name,
        container_name=OUTPUT_BLOB_CONTAINER,
        blob_name=blob_name,
        account_key=account_key,
        permission=BlobSasPermissions(read=True),
        expiry=datetime.now(timezone.utc) + timedelta(hours=sas_hours),
    )

    sas_url = f"{blob.url}?{sas}"
    logger.info(json.dumps({"event": "upload_complete", "url": blob.url, "sas_url_issued": True}))
    return sas_url


@app.function_name(name="TranslatePptx")
@app.route(route="translate/pptx", methods=["POST"], auth_level=AuthLevel.FUNCTION)
def TranslatePptx(req: func.HttpRequest) -> func.HttpResponse:
    logger = logging.getLogger("TranslatePptx")
    job_id = uuid.uuid4().hex
    start = _now_ms()

    # 1) 入力
    try:
        body = req.get_json()
    except Exception:
        return _json(400, {"jobId": job_id, "error": "Invalid JSON"})

    source_url = (body.get("sourceUrl") or "").strip()
    to_lang = (body.get("toLang") or "en").strip()

    if not source_url:
        return _json(400, {"jobId": job_id, "error": "sourceUrl is required"})

    logger.info(json.dumps({"event": "job_start", "jobId": job_id, "toLang": to_lang}))

    in_path = out_path = None
    try:
        # 2) DL（temp保存）
        in_path, in_bytes = download_to_tempfile(source_url, logger)

        # 3) PPTXロード
        t0 = _now_ms()
        prs = Presentation(in_path)
        logger.info(json.dumps({"event": "pptx_load_complete", "ms": _now_ms() - t0}))

        # 4) 抽出→翻訳→反映
        items = extract_items(prs, logger)

        t1 = _now_ms()
        translated = translate_items(items, to_lang, logger)
        logger.info(json.dumps({"event": "translate_complete", "ms": _now_ms() - t1}))

        apply_translations(prs, items, translated, logger)

        # 5) 保存
        out_path = save_to_temp(prs, logger)

        # 6) BlobへUPしてURL返却
        url = upload_blob(out_path, logger)

        return _json(200, {
            "jobId": job_id,
            "status": "ok",
            "inputBytes": in_bytes,
            "textItems": len(items),
            "output": {"type": "blobUrl", "url": url},
            "elapsedMs": _now_ms() - start
        })

    except requests.HTTPError as e:
        logger.exception("HTTP error")
        return _json(502, {"jobId": job_id, "error": "Download/HTTP error", "detail": str(e)})
    except Exception as e:
        logger.exception("Unhandled error")
        return _json(500, {"jobId": job_id, "error": "Unhandled error", "detail": str(e)})
    finally:
        # temp掃除
        for p in [in_path, out_path]:
            if p and os.path.exists(p):
                try:
                    os.remove(p)
                except Exception:
                    pass
